from flask import Flask, render_template, request, send_file, redirect, url_for
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from pathlib import Path
import os
import requests

app = Flask(__name__)

ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
BASE_DIR = Path(__file__).resolve().parent
LOGO_PATH = BASE_DIR / 'iMocha_logo.png'
FIRST_SLIDE_IMAGE = BASE_DIR / 'extracted_images' / 'slide1_img1.jpg'
WEBSITE_TEXT = 'imocha.io'

BRAND_TITLE_COLOR = RGBColor(24, 44, 79)
BRAND_ACCENT_COLOR = RGBColor(247, 129, 33)
BODY_FONT_NAME = 'Poppins'
TITLE_FONT_NAME = 'Playfair Display'
FOOTER_HEIGHT = Inches(1.2)


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_slide_content(slide):
    title = None
    bullets = []
    tables = []
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            table = shape.table
            table_rows = []
            for row in table.rows:
                table_rows.append([cell.text.strip() for cell in row.cells])
            if table_rows:
                tables.append(table_rows)
            continue

        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if not text:
            continue
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not lines:
            continue
        if title is None:
            title = lines[0]
            bullets.extend(lines[1:])
        else:
            bullets.extend(lines)
    return title, bullets, tables


def format_text_frame(text_frame, font_name=BODY_FONT_NAME, font_size=Pt(18), color=BRAND_TITLE_COLOR, bold=False, align=None, line_spacing=None):
    for paragraph in text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        if line_spacing is not None:
            paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color


def get_title_font_size(text):
    if not text:
        return Pt(34)
    length = len(text)
    if length > 100:
        return Pt(22)
    if length > 80:
        return Pt(24)
    if length > 60:
        return Pt(26)
    if length > 45:
        return Pt(28)
    if length > 30:
        return Pt(30)
    return Pt(34)


def get_table_font_size(table_rows):
    max_cell = 0
    cols = len(table_rows[0]) if table_rows else 0
    for row in table_rows:
        for cell_text in row:
            if cell_text and len(cell_text) > max_cell:
                max_cell = len(cell_text)
    if cols >= 5 or max_cell > 80:
        return Pt(10)
    if cols >= 4 or max_cell > 60:
        return Pt(11)
    return Pt(12)


def get_background_image():
    """Fetch and cache a professional background image from the internet."""
    bg_cache_path = BASE_DIR / 'title_background.jpg'
    
    # Use cached image if available
    if bg_cache_path.exists():
        return bg_cache_path
    
    # Professional tech/business background image from Unsplash (free)
    bg_url = 'https://images.unsplash.com/photo-1552664730-d307ca884978?w=1280&q=80'
    
    try:
        response = requests.get(bg_url, timeout=5)
        if response.status_code == 200:
            with open(bg_cache_path, 'wb') as f:
                f.write(response.content)
            return bg_cache_path
    except Exception as e:
        print(f"Background image fetch error: {e}")
    
    return None


def add_slide_footer(slide, prs, include_logo=True):
    if include_logo and LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.4), prs.slide_height - Inches(0.85), width=Inches(1.2))

    footer_box = slide.shapes.add_textbox(prs.slide_width - Inches(3.6), prs.slide_height - Inches(0.75), Inches(3.2), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = WEBSITE_TEXT
    format_text_frame(footer_frame, font_name=BODY_FONT_NAME, font_size=Pt(12), color=BRAND_TITLE_COLOR, bold=False, align=PP_ALIGN.RIGHT)


def create_table_slide(prs, title_text, table_rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.4))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.NONE
    title_frame.text = title_text or 'Table'
    title_frame.paragraphs[0].font.name = TITLE_FONT_NAME
    title_frame.paragraphs[0].font.size = get_title_font_size(title_text)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BRAND_TITLE_COLOR

    rows = len(table_rows)
    cols = len(table_rows[0]) if rows else 0
    if rows and cols:
        table_width = prs.slide_width - Inches(1.0)
        table_top = Inches(1.7)
        table_height = prs.slide_height - table_top - FOOTER_HEIGHT - Inches(0.1)
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), table_top, table_width, table_height)
        table = table_shape.table
        font_size = get_table_font_size(table_rows)
        for col_idx in range(cols):
            table.columns[col_idx].width = int(table_width / cols)

        for r_idx, row_data in enumerate(table_rows):
            for c_idx, cell_text in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = cell_text or ''
                cell_margin = int(Pt(6))
                cell.margin_left = cell_margin
                cell.margin_right = cell_margin
                cell.margin_top = cell_margin
                cell.margin_bottom = cell_margin
                cell.text_frame.word_wrap = True
                cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = Pt(2)
                    for run in paragraph.runs:
                        run.font.name = BODY_FONT_NAME
                        run.font.size = font_size
                        run.font.color.rgb = RGBColor(255, 255, 255) if r_idx == 0 else BRAND_TITLE_COLOR
                        run.font.bold = r_idx == 0
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BRAND_ACCENT_COLOR
                elif r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 244, 250)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_slide_footer(slide, prs)
    return slide


def create_brand_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height

    # Add background image or white background
    bg_image = get_background_image()
    if bg_image:
        try:
            slide.shapes.add_picture(str(bg_image), 0, 0, width=width, height=height)
        except Exception as e:
            print(f"Failed to add background image: {e}")
            background_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
            background_shape.fill.solid()
            background_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            background_shape.line.fill.background()
    else:
        background_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
        background_shape.fill.solid()
        background_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background_shape.line.fill.background()

    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.4), Inches(0.4), width=Inches(1.4))

    # Add white background panel for title readability
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.3), Inches(9.4), Inches(2.8))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bg.line.color.rgb = RGBColor(200, 200, 200)
    title_bg.line.width = Pt(1)

    title_frame = title_bg.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.NONE
    title_frame.margin_left = Pt(12)
    title_frame.margin_right = Pt(12)
    title_frame.margin_top = Pt(14)
    title_frame.margin_bottom = Pt(14)
    title_frame.text = title_text or 'iMocha PPT Formatter'
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.paragraphs[0].font.name = TITLE_FONT_NAME
    title_frame.paragraphs[0].font.size = get_title_font_size(title_text)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BRAND_TITLE_COLOR

    subtitle_paragraph = title_frame.add_paragraph()
    subtitle_paragraph.text = subtitle_text or 'Automatically reformatted and aligned to iMocha brand standards.'
    subtitle_paragraph.font.name = BODY_FONT_NAME
    subtitle_paragraph.font.size = Pt(13)
    subtitle_paragraph.font.color.rgb = BRAND_TITLE_COLOR
    subtitle_paragraph.space_before = Pt(10)

    add_slide_footer(slide, prs, include_logo=False)
    return slide


def create_content_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.4))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.NONE
    title_frame.text = title_text or 'Slide'
    title_frame.paragraphs[0].font.name = TITLE_FONT_NAME
    title_frame.paragraphs[0].font.size = get_title_font_size(title_text)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BRAND_TITLE_COLOR

    body_top = Inches(1.75)
    body_height = prs.slide_height - body_top - FOOTER_HEIGHT - Inches(0.1)
    body_box = slide.shapes.add_textbox(Inches(0.5), body_top, Inches(9.0), body_height)
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    if bullets:
        body_frame.text = '\n'.join(f'• {bullet}' for bullet in bullets)
        format_text_frame(body_frame, font_name=BODY_FONT_NAME, font_size=Pt(14), color=BRAND_TITLE_COLOR, line_spacing=Pt(18))
    else:
        body_frame.text = 'No structured content detected. Please refine or add slide content.'
        format_text_frame(body_frame, font_name=BODY_FONT_NAME, font_size=Pt(14), color=BRAND_TITLE_COLOR, line_spacing=Pt(18))

    add_slide_footer(slide, prs)
    return slide


def build_output_presentation(input_stream):
    incoming = Presentation(input_stream)
    slides_data = []
    first_title = None
    first_subtitle = None

    for idx, slide in enumerate(incoming.slides):
        title, bullets, tables = extract_slide_content(slide)
        if idx == 0:
            first_title = title
            if bullets:
                first_subtitle = ' '.join(bullets[:3])
        slides_data.append((title or f'Slide {idx + 1}', bullets, tables))

    output = Presentation()
    create_brand_title_slide(output, first_title, first_subtitle)

    for title, bullets, tables in slides_data[1:]:
        if bullets:
            create_content_slide(output, title, bullets)
        if tables:
            for table_rows in tables:
                create_table_slide(output, title or 'Table', table_rows)

    output_stream = BytesIO()
    output.save(output_stream)
    output_stream.seek(0)
    return output_stream


@app.route('/')
def index():
    error = request.args.get('error', '')
    return render_template('index.html', error=error)


@app.route('/convert', methods=['POST'])
def convert():
    if 'presentation' not in request.files:
        return redirect(url_for('index', error='No file selected.'))
    file = request.files['presentation']
    if file.filename == '':
        return redirect(url_for('index', error='No file selected.'))
    if not allowed_file(file.filename):
        return redirect(url_for('index', error='Please upload a .ppt or .pptx file.'))

    output_stream = build_output_presentation(file.stream)
    return send_file(
        output_stream,
        as_attachment=True,
        download_name='iMocha_Formatted_Presentation.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


if __name__ == '__main__':
    app.run(debug=True, port=5000)
